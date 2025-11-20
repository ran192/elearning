import os
from tkinter import messagebox, simpledialog
import pandas as pd

# Optional dependencies
try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
    from pptx.util import Inches
except ImportError:
    Presentation = None

def select_file(title="Select file"):
    from tkinter import Tk, filedialog
    root = Tk()
    root.withdraw()
    file_path = filedialog.askopenfilename(
        title=title,
        filetypes=[
            ("Excel files", "*.xlsx *.xls"),
            ("Word files", "*.docx"),
            ("PowerPoint files", "*.pptx")
        ]
    )
    root.destroy()
    return file_path

def save_data(df, original_file_path, custom_name=None):
    if original_file_path is None:
        messagebox.showerror("Error", "Original file path not provided!")
        return

    folder = os.path.dirname(os.path.abspath(original_file_path))
    basename, ext = os.path.splitext(os.path.basename(original_file_path))
    
    if not custom_name:
        from tkinter import Tk
        root = Tk()
        root.withdraw()
        custom_name = simpledialog.askstring("Save As", "Enter file name:", initialvalue=f"{basename}_轉製版")
        root.destroy()
        if not custom_name:
            messagebox.showinfo("Cancelled", "Save cancelled by user.")
            return

    output_path = os.path.join(folder, f"{custom_name}{ext}")

    try:
        if ext in [".xlsx", ".xls"]:
            if df is None:
                messagebox.showerror("Error", "No DataFrame to save!")
                return
            df.to_excel(output_path, index=False)

        elif ext == ".docx":
            if Document is None:
                messagebox.showerror("Error", "python-docx not installed!")
                return
            if df is None:
                messagebox.showerror("Error", "No DataFrame to save!")
                return
            doc = Document()
            doc.add_heading("Data Output", level=1)
            table = doc.add_table(rows=df.shape[0]+1, cols=df.shape[1])
            for j, col_name in enumerate(df.columns):
                table.cell(0, j).text = str(col_name)
            for i in range(df.shape[0]):
                for j in range(df.shape[1]):
                    table.cell(i+1, j).text = str(df.iloc[i,j])
            doc.save(output_path)

        elif ext == ".pptx":
            if Presentation is None:
                messagebox.showerror("Error", "python-pptx not installed!")
                return
            if df is None:
                messagebox.showerror("Error", "No DataFrame to save!")
                return
            prs = Presentation()
            slide_layout = prs.slide_layouts[5]  # blank layout
            slide = prs.slides.add_slide(slide_layout)
            left = top = Inches(1)
            width = Inches(8)
            height = Inches(5)
            textbox = slide.shapes.add_textbox(left, top, width, height)
            tf = textbox.text_frame
            tf.word_wrap = True
            tf.text = df.to_string(index=False)
            prs.save(output_path)
        else:
            messagebox.showerror("Error", f"Unsupported file type: {ext}")
            return

        messagebox.showinfo("Done", f"File saved to:\n{output_path}")
        try:
            os.startfile(output_path)
        except Exception:
            pass
    except Exception as e:
        messagebox.showerror("Error", f"Failed to save file:\n{e}")
